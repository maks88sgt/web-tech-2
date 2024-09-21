from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches

# Path to your PDF file
pdf_path = 'F23 - lesson01.pdf'

# Convert PDF pages to images
images = convert_from_path(pdf_path)

# Create a new PowerPoint presentation
presentation = Presentation()

# Loop through each image (PDF page) and add it as a slide
for image in images:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # Blank slide layout
    # Save each image temporarily
    image_path = 'temp_image.png'
    image.save(image_path, 'PNG')
    # Insert the image into the slide
    slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))

# Save the presentation
pptx_path = 'output_presentation.pptx'
presentation.save(pptx_path)

print(f"Presentation saved to {pptx_path}")
